from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Depends, status, Header
from typing import Optional, List
from fastapi.middleware.cors import CORSMiddleware
from fastapi.security import OAuth2PasswordBearer, OAuth2PasswordRequestForm
from zhipuai import ZhipuAI
from dotenv import load_dotenv
from passlib.context import CryptContext
import jwt
import os
import json
import base64
import fitz # PyMuPDF
from pptx import Presentation
import io
import datetime
import uuid
from sqlalchemy import create_engine, Column, Integer, String, Text, DateTime, ForeignKey, Boolean
from sqlalchemy.ext.declarative import declarative_base
from sqlalchemy.orm import sessionmaker, Session, relationship
from fastapi.staticfiles import StaticFiles

load_dotenv()

app = FastAPI(title="NoteSnap AI API")

# ----------- Auth 设置 -----------
SECRET_KEY = "a_very_secret_key_for_demo_purposes"
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 60 * 24 * 7 # 7 days
pwd_context = CryptContext(schemes=["pbkdf2_sha256"], deprecated="auto")
oauth2_scheme = OAuth2PasswordBearer(tokenUrl="api/login")

# ----------- 数据库设置 -----------
os.makedirs("uploads", exist_ok=True)
app.mount("/uploads", StaticFiles(directory="uploads"), name="uploads")

# 更换新数据库以兼容新增的用户数据表
SQLALCHEMY_DATABASE_URL = "sqlite:///./notes_app_v2.db"
engine = create_engine(SQLALCHEMY_DATABASE_URL, connect_args={"check_same_thread": False})
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)
Base = declarative_base()

class User(Base):
    __tablename__ = "users"
    id = Column(Integer, primary_key=True, index=True)
    username = Column(String, unique=True, index=True)
    email = Column(String, default="")
    phone = Column(String, default="")
    avatar_url = Column(String, default="")
    hashed_password = Column(String)
    created_at = Column(DateTime, default=datetime.datetime.utcnow)
    notes = relationship("NoteRecord", back_populates="owner")

class NoteRecord(Base):
    __tablename__ = "notes"
    id = Column(Integer, primary_key=True, index=True)
    filename = Column(String, index=True)
    title = Column(String, default="Untitled Note")
    category = Column(String, default="Default")
    file_url = Column(String)
    style = Column(String)
    content = Column(Text)
    is_public = Column(Boolean, default=False) # 预留功能：社区分享
    user_id = Column(Integer, ForeignKey("users.id"))
    owner = relationship("User", back_populates="notes")
    created_at = Column(DateTime, default=datetime.datetime.utcnow)

class Comment(Base):
    __tablename__ = "comments"
    id = Column(Integer, primary_key=True, index=True)
    content = Column(String)
    note_id = Column(Integer, ForeignKey("notes.id"))
    user_id = Column(Integer, ForeignKey("users.id"))
    parent_id = Column(Integer, ForeignKey("comments.id"), nullable=True)
    created_at = Column(DateTime, default=datetime.datetime.utcnow)

Base.metadata.create_all(bind=engine)

def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()
# ----------------------------------

# 配置CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

def get_ai_client(custom_api_key: str = None):
    if custom_api_key and "." in custom_api_key:
        return ZhipuAI(api_key=custom_api_key)
    
    load_dotenv(os.path.join(os.path.dirname(__file__), '.env'))
    api_key = os.getenv("ZHIPU_API_KEY", "")
    if not api_key or "." not in api_key:
        raise ValueError("请在 backend/.env 文件中配置 正确的 ZHIPU_API_KEY，或在应用的个人主页中绑定您的本地 API Key")
    return ZhipuAI(api_key=api_key)

# ---------- Auth 路由 ----------
def verify_password(plain_password, hashed_password):
    return pwd_context.verify(plain_password, hashed_password)

def get_password_hash(password):
    return pwd_context.hash(password)

def create_access_token(data: dict):
    to_encode = data.copy()
    expire = datetime.datetime.utcnow() + datetime.timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
    to_encode.update({"exp": expire})
    return jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)

async def get_current_user(token: str = Depends(oauth2_scheme), db: Session = Depends(get_db)):
    credentials_exception = HTTPException(
        status_code=status.HTTP_401_UNAUTHORIZED,
        detail="无效的认证凭据",
        headers={"WWW-Authenticate": "Bearer"},
    )
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        username: str = payload.get("sub")
        if username is None:
            raise credentials_exception
    except jwt.PyJWTError:
        raise credentials_exception
    user = db.query(User).filter(User.username == username).first()
    if user is None:
        raise credentials_exception
    return user

from pydantic import BaseModel
class UserCreate(BaseModel):
    username: str
    password: str
    contact: str = "" # Email or Phone

@app.post("/api/register")
def register_user(user: UserCreate, db: Session = Depends(get_db)):
    db_user = db.query(User).filter(User.username == user.username).first()
    if db_user:
        raise HTTPException(status_code=400, detail="用户名已被注册")
    hashed_pw = get_password_hash(user.password)
    
    is_email = "@" in user.contact
    new_user = User(
        username=user.username, 
        hashed_password=hashed_pw,
        email=user.contact if is_email else "",
        phone=user.contact if not is_email else ""
    )
    db.add(new_user)
    db.commit()
    return {"message": "注册成功"}

@app.post("/api/login")
def login(form_data: OAuth2PasswordRequestForm = Depends(), db: Session = Depends(get_db)):
    user = db.query(User).filter(User.username == form_data.username).first()
    if not user or not verify_password(form_data.password, user.hashed_password):
        raise HTTPException(status_code=400, detail="用户名或密码错误")
    access_token = create_access_token(data={"sub": user.username})
    return {"access_token": access_token, "token_type": "bearer", "username": user.username}

@app.get("/api/me")
def read_users_me(current_user: User = Depends(get_current_user)):
    return {
        "username": current_user.username, 
        "id": current_user.id,
        "email": current_user.email,
        "phone": current_user.phone,
        "avatar_url": getattr(current_user, 'avatar_url', '')
    }

class UserUpdateProfile(BaseModel):
    username: str = None
    email: str = None
    phone: str = None
    avatar_url: str = None

class PasswordChange(BaseModel):
    old_password: str
    new_password: str

@app.post("/api/user/password")
def change_password(data: PasswordChange, db: Session = Depends(get_db), current_user: User = Depends(get_current_user)):
    if not verify_password(data.old_password, current_user.hashed_password):
        raise HTTPException(status_code=400, detail="原密码验证失败，请重试")
    
    current_user.hashed_password = get_password_hash(data.new_password)
    db.commit()
    return {"message": "密码修改成功，请重新登录"}

# ---------- 处理逻辑 ----------
def get_style_prompt(style: str) -> str:
    if "学术论文" in style or "学术" in style:
        return "▶【核心要求】：本材料倾向于学术论文。请严格按照【摘要与背景】、【核心方法】、【实验与数据】、【结论】等学术结构进行梳理，务必保留关键术语、公式与出处的准确性，保持论述严谨严肃。"
    elif "课堂" in style or "考点" in style:
        return "▶【核心要求】：本材料倾向于课堂学习。请将其中的知识点提炼为【期末复习大纲】：包含核心定义、原理及典型例题考法。采用清晰的星号或编号列表展开，极度方便背诵记忆。"
    elif "精简" in style:
        return "▶【核心要求】：请提取最精华、最本质的话语，无需冗长铺垫，砍掉所有边缘性解释，直接把「干货」和「结论」呈现出来，篇幅越少且信息量越大越好。"
    elif not style or style in ["标准", "普通"]:
        return "▶【核心要求】：请均衡提取各章节的核心内容，层级分明，条理清晰，作为一份标准的通用知识笔记。"
    else:
        return f"▶【用户的强制定制风格（最优先应用）】：{style}\n请务必严格遵照本条定制规则进行大纲与内容的摘要构建。"

def process_image(content: bytes, content_type: str, style: str, api_key: str = None):
    client = get_ai_client(api_key)
    base64_image = base64.b64encode(content).decode("utf-8")
    style_guide = get_style_prompt(style)
    
    prompt = f'''你是一个智能笔记整理助手。你需要提取以下图片中的文字内容，并将其归纳为结构化的Markdown笔记。

{style_guide}

【通用排版底线要求】
1. 第一行必须输出以 # 开头的主标题（这是非常重要的一步）；
2. 必须在标题下方单独给出【一句话总结】（加粗标出）；
3. 必须提取【核心关键词】（3-5个，用标签格式，如：#关键词）紧跟在总结下面；
4. 保持原有主体信息准确无误；
5. 为笔记添加合理的章节层级；
6. 对重点内容进行加粗或列表高亮提示；
7. 如果笔记中包含数学公式，请务必使用 $ 包裹行内公式，使用 $$ 包裹块级公式（千万不要使用 \\( \\) 或 \\[ \\]）。'''
    
    response = client.chat.completions.create(
        model="glm-4v",
        messages=[
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:{content_type};base64,{base64_image}"
                        }
                    }
                ]
            }
        ]
    )
    return response.choices[0].message.content

def process_text(text: str, style: str, api_key: str = None):
    client = get_ai_client(api_key)
    max_chars = 100000
    if len(text) > max_chars:
        text = text[:max_chars] + "\n\n...（⚠️文档过长，超维保护触发，此处后续部分已截断）"
    
    style_guide = get_style_prompt(style)

    prompt = f'''你是一个智能笔记整理助手。你需要将以下从文档(PPT/PDF)中提取出的源文本，归纳为结构化的Markdown笔记。

{style_guide}

【通用排版底线要求】
1. 第一行必须输出以 # 开头的主标题（这是最重要的约束）；
2. 必须在标题下方单独给出【一句话总结】（加粗标出）；
3. 必须提取【核心关键词】（3-5个，用标签格式，如：#关键词）紧跟在总结下面；
4. 保持原有主体信息准确无误；
5. 为笔记添加合理的章节层级；
6. 对重点内容进行加粗或列表高亮提示；
7. 如果笔记中包含数学公式，请务必使用 $ 包裹行内公式，使用 $$ 包裹块级公式（千万不要使用 \\( \\) 或 \\[ \\]）。

待整理文本如下：
{text}
'''
    
    response = client.chat.completions.create(
        model="glm-4",
        messages=[
            {
                "role": "user",
                "content": prompt
            }
        ]
    )
    return response.choices[0].message.content

def extract_text_from_pdf(content: bytes):
    text = ""
    with fitz.open(stream=content, filetype="pdf") as doc:
        for page in doc:
            text += page.get_text() + "\n"
    return text

def extract_text_from_pptx(content: bytes):
    prs = Presentation(io.BytesIO(content))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

@app.post("/api/process-note")
async def process_note(
    files: List[UploadFile] = File(...), 
    style: str = Form("标准"), 
    category: str = Form("未分类"),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
    x_zhipu_api_key: str = Header(None)
):
    try:
        accumulated_text = ""
        base64_images = []
        file_urls = []
        file_names = []
        
        for file in files:
            filename = file.filename.lower()
            file_names.append(file.filename)
            content = await file.read()
            
            # 1. 保存原始文件
            unique_filename = f"{uuid.uuid4()}_{file.filename}"
            file_path = os.path.join("uploads", unique_filename)
            with open(file_path, "wb") as f:
                f.write(content)
            file_urls.append(f"http://localhost:8000/uploads/{unique_filename}")
            
            # 2. 提取内容
            if file.content_type.startswith("image/"):
                b64 = base64.b64encode(content).decode("utf-8")
                base64_images.append((b64, file.content_type))
            elif filename.endswith(".pdf") or file.content_type == "application/pdf":
                extracted = extract_text_from_pdf(content)
                if extracted.strip():
                    accumulated_text += f"\n--- 文档：{file.filename} ---\n{extracted}"
            elif filename.endswith(".pptx") or "presentation" in file.content_type:
                extracted = extract_text_from_pptx(content)
                if extracted.strip():
                    accumulated_text += f"\n--- 演示文稿：{file.filename} ---\n{extracted}"
            else:
                raise HTTPException(status_code=400, detail=f"文件 {file.filename} 不支持，仅支持图片、PDF或PPTX。")
        
        # 统一去让大模型处理
        client = get_ai_client(x_zhipu_api_key)
        style_guide = get_style_prompt(style)
        
        if len(base64_images) > 0:
            prompt = f'''你是一个智能笔记整理助手。你需要提取以下用户上传的图片与文档等混合内容，并将其归纳为结构化的Markdown笔记。
{style_guide}
【通用排版底线要求】
1. 第一行首出以 # 开头的主标题
2. 标题下方给出【一句话总结】（加粗标出）
3. 提取【核心关键词】（3-5个，#关键词 格式）
4. 其他重点加粗、分类清晰、数学公式用 $ / $$ 包裹。
以下是提取出来的额外文档文本（如果没有可以忽略图片即可）：
{accumulated_text}'''
            
            msg_content = [{"type": "text", "text": prompt}]
            for b64, ctype in base64_images:
                msg_content.append({"type": "image_url", "image_url": {"url": f"data:{ctype};base64,{b64}"}})
            
            response = client.chat.completions.create(
                model="glm-4v",
                messages=[{"role": "user", "content": msg_content}]
            )
            result = response.choices[0].message.content
        else:
            if not accumulated_text.strip():
                raise ValueError("未能提取到任何有效文本或图片。")
            max_chars = 100000
            if len(accumulated_text) > max_chars:
                accumulated_text = accumulated_text[:max_chars] + "\n\n...（⚠️文档过长，此处后续部分已截断）"
            
            prompt = f'''你是一个智能笔记整理助手。你需要将以下源文本，归纳为结构化的Markdown笔记。
{style_guide}
【排版底线要求】同上。
待整理文本如下：
{accumulated_text}'''

            response = client.chat.completions.create(
                model="glm-4",
                messages=[{"role": "user", "content": prompt}]
            )
            result = response.choices[0].message.content
            
        # 清理大模型可能带有的 ```markdown 与 ``` 的外包装
        if result.strip().startswith("```markdown"):
            result = result.strip()[11:]
        elif result.strip().startswith("```"):
            result = result.strip()[3:]
        if result.endswith("```"):
            result = result[:-3]
        result = result.strip()
            
        # 自动提取Markdown中第一个真实标题作为Note的标题
        import re
        title_match = None
        for match in re.finditer(r'^[\*\s]*#{1,2}\s*[\*\s]*(.*?)(?=[\*\s]*(\n|$))', result, re.MULTILINE):
            extracted = match.group(1).strip()
            if "总结" not in extracted and "关键词" not in extracted and len(extracted) > 1:
                title_match = match
                break

        if title_match:
            generated_title = title_match.group(1).strip()
            result = result.replace(title_match.group(0), "", 1).strip()
        else:
            generated_title = "、".join(file_names)[:50]

        # 3. 存入当前用户数据库
        file_url_str = ",".join(file_urls)
        db_note = NoteRecord(filename=",".join(file_names), title=generated_title, file_url=file_url_str, style=style, category=category, content=result, user_id=current_user.id)
        db.add(db_note)
        db.commit()
        db.refresh(db_note)

        return {"note": result, "id": db_note.id, "title": db_note.title, "category": db_note.category, "file_url": file_url_str}

    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/notes")
def get_notes(
    skip: int = 0, limit: int = 100, 
    search: str = None,
    db: Session = Depends(get_db), 
    current_user: User = Depends(get_current_user)
):
    query = db.query(NoteRecord).filter(NoteRecord.user_id == current_user.id)
    if search:
        query = query.filter(
            (NoteRecord.title.ilike(f"%{search}%")) |
            (NoteRecord.content.ilike(f"%{search}%"))
        )
    notes = query.order_by(NoteRecord.created_at.desc()).offset(skip).limit(limit).all()
    result = []
    for n in notes:
        result.append({
            "id": n.id,
            "filename": n.filename,
            "title": n.title,
            "category": n.category,
            "created_at": n.created_at.isoformat(),
            "content": n.content,
            "file_url": n.file_url,
            "style": n.style,
            "is_public": n.is_public
        })
    return result

@app.get("/api/community/notes")
def get_community_notes(search: str = None, db: Session = Depends(get_db)):
    # 获取分享的笔记
    query = db.query(NoteRecord).filter(NoteRecord.is_public == True)
    if search:
        query = query.filter(
            (NoteRecord.title.ilike(f"%{search}%")) |
            (NoteRecord.content.ilike(f"%{search}%"))
        )
    notes = query.order_by(NoteRecord.created_at.desc()).limit(50).all()
    result = []
    for n in notes:
        result.append({
            "id": n.id,
            "filename": n.filename,
            "title": n.title,
            "category": n.category,
            "created_at": n.created_at.isoformat(),
            "content": n.content[:200] + "...", # 社区只返回摘要
            "style": n.style,
            "author": n.owner.username if n.owner else "匿名用户"
        })
    return result

@app.get("/api/community/notes/{note_id}")
def get_community_note_detail(note_id: int, db: Session = Depends(get_db)):
    # 允许社区内获取单篇公开笔记全文
    note = db.query(NoteRecord).filter(NoteRecord.id == note_id, NoteRecord.is_public == True).first()
    if not note:
        raise HTTPException(status_code=404, detail="未找到该分享笔记信息。")
    return {
        "id": note.id,
        "filename": note.filename,
        "title": note.title,
        "category": note.category,
        "created_at": note.created_at.isoformat(),
        "content": note.content,
        "style": note.style,
        "author": note.owner.username if note.owner else "匿名用户"
    }

@app.patch("/api/notes/{note_id}/public")
def toggle_public(note_id: int, is_public: bool, db: Session = Depends(get_db), current_user: User = Depends(get_current_user)):
    note = db.query(NoteRecord).filter(NoteRecord.id == note_id, NoteRecord.user_id == current_user.id).first()
    if not note:
        raise HTTPException(status_code=404, detail="笔记找不到了或无操作权限")
    note.is_public = is_public
    db.commit()
    return {"message": "鐘舵€佹洿鏂版垚鍔?", "is_public": note.is_public}


@app.post("/api/notes/{note_id}/quiz")
def generate_quiz(
    note_id: int, 
    db: Session = Depends(get_db), 
    current_user: User = Depends(get_current_user),
    x_zhipu_api_key: str = Header(None)
):
    note = db.query(NoteRecord).filter(NoteRecord.id == note_id).first()
    if not note:
        raise HTTPException(status_code=404, detail="笔记不存在")
    
    if not note.is_public and note.user_id != current_user.id:
        raise HTTPException(status_code=403, detail="没有权限访问该笔记")

    client = get_ai_client(x_zhipu_api_key)
    
    prompt = f"""你是一个擅长提炼知识点的AI助教。请仔细阅读下面的笔记内容，并生成3道单项选择题。
必须严格输出纯 JSON 数组格式（不要带有任何 ```json 或者其它说明语言）。
格式要求如下：
[
  {{
    "question": "问题是什么？",
    "options": ["选项A的内容", "选项B的内容", "选项C的内容", "选项D的内容"],
    "answer": 0,
    "explanation": "为什么选A的简短解析"
  }}
]

【笔记内容】：
{note.content}
"""

    try:
        response = client.chat.completions.create(
            model="glm-4",
            messages=[
                {"role": "system", "content": "你是一个只会输出合法JSON数组的API，不输出任何其他字。"},
                {"role": "user", "content": prompt}
            ]
        )
        answer = response.choices[0].message.content.strip()
        if answer.startswith("```json"): answer = answer[7:]
        elif answer.startswith("```"): answer = answer[3:]
        if answer.endswith("```"): answer = answer[:-3]
        
        quiz_data = json.loads(answer.strip())
        return quiz_data
    except Exception as e:
        print(f"Quiz Error: {str(e)}\nAI output was: {answer}")
        raise HTTPException(status_code=500, detail="生成测验失败，笔记内容可能不足或格式错误")


class NoteChatRequest(BaseModel):
    message: str

@app.post("/api/notes/{note_id}/chat")
def chat_with_note(
    note_id: int, 
    chat_request: NoteChatRequest, 
    db: Session = Depends(get_db), 
    current_user: User = Depends(get_current_user),
    x_zhipu_api_key: str = Header(None)
):
    note = db.query(NoteRecord).filter(NoteRecord.id == note_id).first()
    if not note:
        raise HTTPException(status_code=404, detail="笔记不存在")
    
    if not note.is_public and note.user_id != current_user.id:
        raise HTTPException(status_code=403, detail="没有权限访问该笔记")

    client = get_ai_client(x_zhipu_api_key)
    
    prompt = f"""你是一个智能笔记助手。以下是用户的一篇笔记（可能包含代码或文字）。
请基于提供的笔记内容，回答用户的问题。如果用户的提问超出了笔记内容，可以使用你的知识进行扩展，但请优先参考笔记。
如果是代码，请在必要时以Markdown格式输出。

【笔记标题】: {note.title}
【笔记分类】: {note.category}
【笔记内容】:
{note.content}

用户的问题是: {chat_request.message}"""

    try:
        response = client.chat.completions.create(
            model="glm-4",
            messages=[
                {"role": "system", "content": "你是一个基于用户笔记内容的智能助手，请主要根据提供的笔记内容回答，语气要专业友好。"},
                {"role": "user", "content": prompt}
            ]
        )
        answer = response.choices[0].message.content
        return {"answer": answer}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

class CommentCreate(BaseModel):
    content: str
    parent_id: Optional[int] = None

@app.post("/api/notes/{note_id}/comments")
def add_comment(note_id: int, comment: CommentCreate, db: Session = Depends(get_db), current_user: User = Depends(get_current_user)):
    db_note = db.query(NoteRecord).filter(NoteRecord.id == note_id).first()
    if not db_note:
        raise HTTPException(status_code=404, detail="绗旇涓嶅瓨鍦¨")
    
    new_comment = Comment(
        content=comment.content,
        note_id=note_id,
        user_id=current_user.id,
        parent_id=comment.parent_id
    )
    db.add(new_comment)
    db.commit()
    db.refresh(new_comment)
    
    return {
        "id": new_comment.id,
        "content": new_comment.content,
        "created_at": new_comment.created_at.isoformat(),
        "author": current_user.username,
        "author_avatar": getattr(current_user, 'avatar_url', ''),
        "parent_id": new_comment.parent_id
    }

@app.get("/api/notes/{note_id}/comments")
def get_comments(note_id: int, db: Session = Depends(get_db)):
    db_note = db.query(NoteRecord).filter(NoteRecord.id == note_id).first()
    if not db_note:
        raise HTTPException(status_code=404, detail="绗旇涓嶅瓨鍦¨")
        
    comments = db.query(Comment).filter(Comment.note_id == note_id).order_by(Comment.created_at.asc()).all()
    result = []
    
    # Pre-fetch user dictionary for performance and easy reply_to mapping
    users_dict = {u.id: {"username": u.username, "avatar_url": getattr(u, 'avatar_url', '')} for u in db.query(User).all()}
    
    for c in comments:
        user_info = users_dict.get(c.user_id, {"username": "匿名", "avatar_url": ""})
        result.append({
            "id": c.id,
            "content": c.content,
            "created_at": c.created_at.isoformat(),
            "author": user_info["username"],
            "author_avatar": user_info["avatar_url"],
            "parent_id": c.parent_id
        })
    return result


@app.delete("/api/notes/{note_id}")
def delete_note(note_id: int, db: Session = Depends(get_db), current_user: User = Depends(get_current_user)):
    db_note = db.query(NoteRecord).filter(NoteRecord.id == note_id, NoteRecord.user_id == current_user.id).first()
    if not db_note:
        raise HTTPException(status_code=404, detail="笔记不存在或无权限删除")
    
    db.delete(db_note)
    db.commit()
    return {"message": "笔记已删除"}

class NoteCreate(BaseModel):
    title: str = "Untitled Note"
    content: str
    category: str = "未分类"
    is_public: bool = False

@app.post("/api/notes")
def create_raw_note(note: NoteCreate, db: Session = Depends(get_db), current_user: User = Depends(get_current_user)):
    db_note = NoteRecord(
        title=note.title,
        content=note.content,
        category=note.category,
        is_public=note.is_public,
        user_id=current_user.id,
        filename="Manual Note",
        style="standard"
    )
    db.add(db_note)
    db.commit()
    db.refresh(db_note)
    return {
        "id": db_note.id, 
        "title": db_note.title, 
        "category": db_note.category, 
        "message": "笔记创建成功"
    }


class NoteUpdate(BaseModel):
    title: str = None
    content: str = None
    category: str = None

@app.put("/api/notes/{note_id}")
def update_note(note_id: int, note_update: NoteUpdate, db: Session = Depends(get_db), current_user: User = Depends(get_current_user)):
    db_note = db.query(NoteRecord).filter(NoteRecord.id == note_id, NoteRecord.user_id == current_user.id).first()
    if not db_note:
        raise HTTPException(status_code=404, detail="笔记找不到了或无操作权限")
    
    if note_update.title is not None:
        db_note.title = note_update.title
    if note_update.content is not None:
        db_note.content = note_update.content
    if note_update.category is not None:
        db_note.category = note_update.category
        
    db.commit()
    return {"message": "笔记已保存"}


@app.put("/api/me")
def update_user_profile(profile_data: UserUpdateProfile, db: Session = Depends(get_db), current_user: User = Depends(get_current_user)):
    if profile_data.username and profile_data.username != current_user.username:
        existing = db.query(User).filter(User.username == profile_data.username).first()
        if existing:
            raise HTTPException(status_code=400, detail="鐢ㄦ埛鍚嶅凡瀛樺湪")
        current_user.username = profile_data.username
    if profile_data.email is not None:
        current_user.email = profile_data.email
    if profile_data.phone is not None:
        current_user.phone = profile_data.phone
    if profile_data.avatar_url is not None:
        current_user.avatar_url = profile_data.avatar_url
    
    db.commit()
    db.refresh(current_user)
    return {
        "username": current_user.username, 
        "id": current_user.id,
        "email": current_user.email,
        "phone": current_user.phone,
        "avatar_url": getattr(current_user, 'avatar_url', '')
    }

@app.post("/api/user/avatar")
async def upload_avatar(
    file: UploadFile = File(...),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    try:
        content = await file.read()
        import uuid
        import os
        import json
        unique_filename = f"avatar_{uuid.uuid4()}_{file.filename}"
        file_path = os.path.join("uploads", unique_filename)
        with open(file_path, "wb") as f:
            f.write(content)
        file_url = f"http://localhost:8000/uploads/{unique_filename}"
        
        current_user.avatar_url = file_url
        db.commit()
        return {"avatar_url": file_url}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))







if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
